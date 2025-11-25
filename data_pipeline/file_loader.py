import hashlib
import logging
import mimetypes
import os
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

logger = logging.getLogger(__name__)


@dataclass
class FileLoaderConfig:
    """Configuration for file loader."""
    supported_extensions: List[str] = field(default_factory=lambda: [
        ".pdf", ".docx", ".doc", ".txt", ".md", ".html",
        ".csv", ".xlsx", ".xls", ".pptx", ".ppt",
        ".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tiff"
    ])
    max_file_size_mb: float = 100.0
    encoding: str = "utf-8"
    extract_images: bool = True
    extract_tables: bool = True


@dataclass
class LoadedDocument:
    """Represents a loaded and parsed document."""
    file_path: str
    file_name: str
    file_extension: str
    content: str
    metadata: Dict[str, Any]
    file_hash: str
    loaded_at: str
    images: List[Dict[str, Any]] = field(default_factory=list)
    tables: List[Dict[str, Any]] = field(default_factory=list)
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dictionary."""
        return {
            "file_path": self.file_path,
            "file_name": self.file_name,
            "file_extension": self.file_extension,
            "content": self.content,
            "metadata": self.metadata,
            "file_hash": self.file_hash,
            "loaded_at": self.loaded_at,
            "images": self.images,
            "tables": self.tables,
        }


class FileLoaderError(Exception):
    """Exception raised for file loading errors."""
    
    def __init__(self, message: str, file_path: str = "", original_error: Optional[Exception] = None):
        self.message = message
        self.file_path = file_path
        self.original_error = original_error
        super().__init__(f"{message} (file: {file_path})")


class FileLoader:
    """
    Universal file loader supporting multiple document formats.
    
    Supported formats:
    - PDF (with OCR support for scanned documents)
    - Office documents (docx, doc, xlsx, xls, pptx, ppt)
    - Text files (txt, md)
    - HTML
    - CSV
    - Images (for OCR processing)
    """
    
    def __init__(self, config: Optional[FileLoaderConfig] = None):
        """
        Initialize the file loader.
        
        Args:
            config: File loader configuration
        """
        self.config = config or FileLoaderConfig()
        self._loaders = {}
        self._initialize_loaders()
    
    def _initialize_loaders(self) -> None:
        """Initialize format-specific loaders."""
        # Register loaders for each format
        self._loaders = {
            ".pdf": self._load_pdf,
            ".docx": self._load_docx,
            ".doc": self._load_doc,
            ".txt": self._load_text,
            ".md": self._load_text,
            ".html": self._load_html,
            ".htm": self._load_html,
            ".csv": self._load_csv,
            ".xlsx": self._load_excel,
            ".xls": self._load_excel,
            ".pptx": self._load_pptx,
            ".ppt": self._load_pptx,
            ".png": self._load_image,
            ".jpg": self._load_image,
            ".jpeg": self._load_image,
            ".gif": self._load_image,
            ".bmp": self._load_image,
            ".tiff": self._load_image,
        }
    
    def _calculate_file_hash(self, file_path: str) -> str:
        """Calculate MD5 hash of file."""
        hash_md5 = hashlib.md5()
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(4096), b""):
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
    
    def _get_file_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract file metadata."""
        path = Path(file_path)
        stat = path.stat()
        
        return {
            "file_name": path.name,
            "file_extension": path.suffix.lower(),
            "file_size_bytes": stat.st_size,
            "file_size_mb": round(stat.st_size / (1024 * 1024), 2),
            "created_time": datetime.fromtimestamp(stat.st_ctime).isoformat(),
            "modified_time": datetime.fromtimestamp(stat.st_mtime).isoformat(),
            "mime_type": mimetypes.guess_type(file_path)[0],
        }
    
    def _validate_file(self, file_path: str) -> None:
        """Validate file before loading."""
        path = Path(file_path)
        
        if not path.exists():
            raise FileLoaderError(f"File not found", file_path)
        
        if not path.is_file():
            raise FileLoaderError(f"Path is not a file", file_path)
        
        extension = path.suffix.lower()
        if extension not in self.config.supported_extensions:
            raise FileLoaderError(
                f"Unsupported file extension: {extension}. "
                f"Supported: {self.config.supported_extensions}",
                file_path
            )
        
        file_size_mb = path.stat().st_size / (1024 * 1024)
        if file_size_mb > self.config.max_file_size_mb:
            raise FileLoaderError(
                f"File too large: {file_size_mb:.2f}MB > {self.config.max_file_size_mb}MB",
                file_path
            )
    
    def _load_pdf(self, file_path: str) -> tuple:
        """Load PDF document."""
        content = ""
        images = []
        tables = []
        
        try:
            # Try PyMuPDF first
            import fitz
            
            doc = fitz.open(file_path)
            text_parts = []
            
            for page_num, page in enumerate(doc):
                text = page.get_text()
                text_parts.append(text)
                
                # Extract images if enabled
                if self.config.extract_images:
                    image_list = page.get_images()
                    for img_index, img in enumerate(image_list):
                        images.append({
                            "page": page_num + 1,
                            "index": img_index,
                            "xref": img[0],
                        })
            
            content = "\n\n".join(text_parts)
            doc.close()
            
        except ImportError:
            # Fallback to pypdf
            try:
                from pypdf import PdfReader
                
                reader = PdfReader(file_path)
                text_parts = []
                
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                
                content = "\n\n".join(text_parts)
                
            except ImportError:
                raise FileLoaderError(
                    "No PDF library available. Install PyMuPDF or pypdf",
                    file_path
                )
        
        return content, images, tables
    
    def _load_docx(self, file_path: str) -> tuple:
        """Load DOCX document."""
        try:
            from docx import Document
            
            doc = Document(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = "\n\n".join(paragraphs)
            
            # Extract tables
            tables = []
            if self.config.extract_tables:
                for table_idx, table in enumerate(doc.tables):
                    table_data = []
                    for row in table.rows:
                        row_data = [cell.text for cell in row.cells]
                        table_data.append(row_data)
                    tables.append({
                        "index": table_idx,
                        "data": table_data
                    })
            
            return content, [], tables
            
        except ImportError:
            raise FileLoaderError(
                "python-docx is required. Install with: pip install python-docx",
                file_path
            )
    
    def _load_doc(self, file_path: str) -> tuple:
        """Load DOC document (legacy format)."""
        try:
            # Try textract for .doc files
            import textract  # pyright: ignore[reportMissingImports]
            
            content = textract.process(file_path).decode(self.config.encoding)
            return content, [], []
            
        except ImportError:
            # Fallback: try to read as text
            logger.warning(
                f"textract not available for .doc files. "
                f"Install with: pip install textract"
            )
            return self._load_text(file_path)
        except Exception as e:
            raise FileLoaderError(f"Failed to load .doc file: {str(e)}", file_path, e)
    
    def _load_text(self, file_path: str) -> tuple:
        """Load plain text file."""
        encodings = [self.config.encoding, "utf-8", "latin-1", "cp1252"]
        
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()
                return content, [], []
            except UnicodeDecodeError:
                continue
        
        raise FileLoaderError(f"Could not decode file with any encoding", file_path)
    
    def _load_html(self, file_path: str) -> tuple:
        """Load HTML document."""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, "r", encoding=self.config.encoding) as f:
                soup = BeautifulSoup(f.read(), "html.parser")
            
            # Remove script and style elements
            for element in soup(["script", "style", "nav", "footer", "header"]):
                element.decompose()
            
            content = soup.get_text(separator="\n", strip=True)
            
            # Extract tables
            tables = []
            if self.config.extract_tables:
                for table_idx, table in enumerate(soup.find_all("table")):
                    table_data = []
                    for row in table.find_all("tr"):
                        row_data = [cell.get_text(strip=True) for cell in row.find_all(["td", "th"])]
                        table_data.append(row_data)
                    tables.append({
                        "index": table_idx,
                        "data": table_data
                    })
            
            return content, [], tables
            
        except ImportError:
            # Fallback: basic text extraction
            with open(file_path, "r", encoding=self.config.encoding) as f:
                content = f.read()
            # Simple HTML tag removal
            import re
            content = re.sub(r"<[^>]+>", " ", content)
            return content, [], []
    
    def _load_csv(self, file_path: str) -> tuple:
        """Load CSV file."""
        import csv
        
        with open(file_path, "r", encoding=self.config.encoding, newline="") as f:
            reader = csv.reader(f)
            rows = list(reader)
        
        # Convert to text representation
        content_parts = []
        for row in rows:
            content_parts.append(" | ".join(row))
        
        content = "\n".join(content_parts)
        tables = [{"index": 0, "data": rows}]
        
        return content, [], tables
    
    def _load_excel(self, file_path: str) -> tuple:
        """Load Excel file."""
        try:
            import openpyxl
            
            wb = openpyxl.load_workbook(file_path, data_only=True)
            content_parts = []
            tables = []
            
            for sheet_idx, sheet_name in enumerate(wb.sheetnames):
                sheet = wb[sheet_name]
                content_parts.append(f"## Sheet: {sheet_name}")
                
                sheet_data = []
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) if cell is not None else "" for cell in row]
                    if any(row_values):  # Skip empty rows
                        content_parts.append(" | ".join(row_values))
                        sheet_data.append(row_values)
                
                tables.append({
                    "sheet": sheet_name,
                    "index": sheet_idx,
                    "data": sheet_data
                })
                
                content_parts.append("")  # Empty line between sheets
            
            content = "\n".join(content_parts)
            return content, [], tables
            
        except ImportError:
            raise FileLoaderError(
                "openpyxl is required for Excel files. Install with: pip install openpyxl",
                file_path
            )
    
    def _load_pptx(self, file_path: str) -> tuple:
        """Load PowerPoint file."""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content_parts = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts.append(f"## Slide {slide_num}")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)
                
                content_parts.append("")
            
            content = "\n".join(content_parts)
            return content, [], []
            
        except ImportError:
            raise FileLoaderError(
                "python-pptx is required. Install with: pip install python-pptx",
                file_path
            )
    
    def _load_image(self, file_path: str) -> tuple:
        """Load image file (returns metadata, actual OCR done separately)."""
        try:
            from PIL import Image
            
            with Image.open(file_path) as img:
                metadata = {
                    "width": img.width,
                    "height": img.height,
                    "format": img.format,
                    "mode": img.mode,
                }
            
            # Return placeholder - actual OCR will be done by OCR processor
            content = f"[Image file: {os.path.basename(file_path)}]"
            images = [{
                "path": file_path,
                "metadata": metadata,
                "requires_ocr": True
            }]
            
            return content, images, []
            
        except ImportError:
            raise FileLoaderError(
                "Pillow is required for images. Install with: pip install Pillow",
                file_path
            )
    
    def load(self, file_path: str) -> LoadedDocument:
        """
        Load a document from file.
        
        Args:
            file_path: Path to the document
            
        Returns:
            LoadedDocument with parsed content
        """
        file_path = str(Path(file_path).resolve())
        
        # Validate file
        self._validate_file(file_path)
        
        # Get metadata
        metadata = self._get_file_metadata(file_path)
        file_hash = self._calculate_file_hash(file_path)
        
        # Get extension
        extension = Path(file_path).suffix.lower()
        
        # Load using appropriate loader
        loader = self._loaders.get(extension)
        if not loader:
            raise FileLoaderError(f"No loader available for {extension}", file_path)
        
        try:
            content, images, tables = loader(file_path)
            
            # Clean content
            content = self._clean_content(content)
            
            logger.info(f"Loaded document: {file_path} ({len(content)} chars)")
            
            return LoadedDocument(
                file_path=file_path,
                file_name=metadata["file_name"],
                file_extension=extension,
                content=content,
                metadata=metadata,
                file_hash=file_hash,
                loaded_at=datetime.now().isoformat(),
                images=images,
                tables=tables,
            )
            
        except FileLoaderError:
            raise
        except Exception as e:
            raise FileLoaderError(f"Failed to load document: {str(e)}", file_path, e)
    
    def _clean_content(self, content: str) -> str:
        """Clean and normalize content."""
        if not content:
            return ""
        
        # Normalize whitespace
        import re
        
        # Replace multiple spaces with single space
        content = re.sub(r"[ \t]+", " ", content)
        
        # Replace multiple newlines with double newline
        content = re.sub(r"\n{3,}", "\n\n", content)
        
        # Strip leading/trailing whitespace
        content = content.strip()
        
        return content
    
    def load_directory(
        self,
        directory: str,
        recursive: bool = True,
        extensions: Optional[List[str]] = None
    ) -> List[LoadedDocument]:
        """
        Load all documents from a directory.
        
        Args:
            directory: Directory path
            recursive: Whether to search subdirectories
            extensions: Filter by extensions (uses config if None)
            
        Returns:
            List of LoadedDocument objects
        """
        directory = Path(directory)
        
        if not directory.exists():
            raise FileLoaderError(f"Directory not found: {directory}", str(directory))
        
        if not directory.is_dir():
            raise FileLoaderError(f"Path is not a directory: {directory}", str(directory))
        
        extensions = extensions or self.config.supported_extensions
        documents = []
        errors = []
        
        # Get all files
        if recursive:
            files = list(directory.rglob("*"))
        else:
            files = list(directory.glob("*"))
        
        # Filter and load
        for file_path in files:
            if file_path.is_file() and file_path.suffix.lower() in extensions:
                try:
                    doc = self.load(str(file_path))
                    documents.append(doc)
                except FileLoaderError as e:
                    logger.warning(f"Failed to load {file_path}: {e.message}")
                    errors.append({"file": str(file_path), "error": e.message})
        
        logger.info(
            f"Loaded {len(documents)} documents from {directory} "
            f"({len(errors)} errors)"
        )
        
        return documents
    
    def get_supported_extensions(self) -> List[str]:
        """Get list of supported file extensions."""
        return list(self._loaders.keys())

